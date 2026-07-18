from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parents[1] / "deliverables"
OUT_FILE = OUT_DIR / "Codex_Go_项目宣传介绍.pptx"

W, H = 13.333, 7.5

# Palette
BG = "07111F"
BG2 = "0C1A2B"
PANEL = "10243A"
PANEL2 = "152D46"
WHITE = "F6FAFF"
TEXT = "D8E6F3"
MUTED = "8EA5BA"
CYAN = "46D9FF"
GREEN = "65E6A5"
LIME = "B8F36B"
BLUE = "4B7BFF"
ORANGE = "FFB15A"
RED = "FF6B7A"
GRID = "203750"
FONT = "Microsoft YaHei"
MONO = "Cascadia Mono"


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color, width=1, transparency=0):
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=True, transparency=0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shp, fill, transparency)
    if line:
        set_line(shp, line)
    else:
        shp.line.fill.background()
    return shp


def circle(slide, x, y, d, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    set_fill(shp, fill)
    if line:
        set_line(shp, line)
    else:
        shp.line.fill.background()
    return shp


def line(slide, x1, y1, x2, y2, color=GRID, width=1.2, dash=None):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    set_line(shp, color, width)
    if dash:
        shp.line.dash_style = dash
    return shp


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    size=18,
    color=TEXT,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def rich_text(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for item in runs:
        run = p.add_run()
        run.text = item["text"]
        run.font.name = item.get("font", FONT)
        run.font.size = Pt(item.get("size", size))
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = rgb(item.get("color", TEXT))
    return box


def add_bg(slide, variant=0):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(BG if variant == 0 else BG2)
    # Subtle technical grid.
    for x in [0.7, 3.7, 6.7, 9.7, 12.7]:
        line(slide, x, 0, x, H, GRID, 0.45)
    for y in [1.0, 3.0, 5.0, 7.0]:
        line(slide, 0, y, W, y, GRID, 0.45)
    # Accent rails.
    rect(slide, 0, 0, 0.08, H, CYAN, radius=False)
    rect(slide, 0.08, 0, 0.035, 2.1, LIME, radius=False)


def add_header(slide, index, title_value, kicker=None):
    if kicker:
        text(slide, kicker.upper(), 0.62, 0.37, 7.0, 0.28, 10, CYAN, True)
    text(slide, title_value, 0.62, 0.68, 11.6, 0.64, 28, WHITE, True)
    text(slide, f"{index:02d}", 12.25, 0.46, 0.48, 0.34, 11, MUTED, True, align=PP_ALIGN.RIGHT)
    line(slide, 0.62, 1.38, 12.72, 1.38, GRID, 1)


def footer(slide, note="Codex Go · Project Introduction"):
    text(slide, note, 0.62, 7.14, 6.5, 0.2, 8, MUTED)
    text(slide, "2026.07", 11.8, 7.14, 0.9, 0.2, 8, MUTED, align=PP_ALIGN.RIGHT)


def tag(slide, value, x, y, w, color=CYAN):
    rect(slide, x, y, w, 0.34, fill=PANEL2, line=color, radius=True)
    text(slide, value, x, y + 0.015, w, 0.25, 10, color, True, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, title_value, body, accent=CYAN, badge=None):
    rect(slide, x, y, w, h, PANEL, GRID)
    rect(slide, x, y, 0.06, h, accent, radius=False)
    if badge:
        circle(slide, x + 0.28, y + 0.25, 0.46, accent)
        text(slide, badge, x + 0.28, y + 0.26, 0.46, 0.35, 15, BG, True, align=PP_ALIGN.CENTER)
        tx = x + 0.88
        tw = w - 1.12
    else:
        tx = x + 0.28
        tw = w - 0.54
    text(slide, title_value, tx, y + 0.24, tw, 0.36, 16, WHITE, True)
    text(slide, body, tx, y + 0.72, tw, h - 0.9, 11, TEXT, False, line_spacing=1.13)


def pill(slide, value, x, y, w, accent=GREEN):
    rect(slide, x, y, w, 0.36, PANEL2, accent)
    text(slide, value, x, y + 0.035, w, 0.23, 9, accent, True, align=PP_ALIGN.CENTER)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    # Decorative orbit / agent core.
    circle(slide, 8.3, 1.05, 4.1, PANEL)
    circle(slide, 8.95, 1.70, 2.8, BG2, CYAN)
    circle(slide, 9.77, 2.52, 1.16, LIME)
    text(slide, "GO", 9.77, 2.81, 1.16, 0.4, 23, BG, True, align=PP_ALIGN.CENTER)
    for x, y, c, label in [
        (8.07, 1.15, CYAN, "MCP"),
        (11.62, 1.67, GREEN, "TUI"),
        (8.18, 4.38, ORANGE, "CLI"),
        (11.35, 4.69, BLUE, "API"),
    ]:
        circle(slide, x, y, 0.65, c)
        text(slide, label, x, y + 0.19, 0.65, 0.22, 9, BG, True, align=PP_ALIGN.CENTER)
        line(slide, x + 0.33, y + 0.33, 10.35, 3.10, c, 1.4)
    tag(slide, "AI CODING AGENT", 0.72, 0.72, 2.05, CYAN)
    text(slide, "Codex Go", 0.72, 1.56, 7.0, 0.9, 43, WHITE, True)
    text(slide, "用 Go 重构的新一代\n跨平台 AI 编程代理", 0.72, 2.55, 7.0, 1.45, 25, TEXT, True, line_spacing=1.08)
    text(
        slide,
        "更易部署 · 更易集成 · 更可治理",
        0.75,
        4.32,
        6.1,
        0.45,
        17,
        LIME,
        True,
    )
    pill(slide, "Interactive TUI", 0.75, 5.13, 1.75, CYAN)
    pill(slide, "Headless Exec", 2.67, 5.13, 1.66, GREEN)
    pill(slide, "MCP / Plugin", 4.50, 5.13, 1.55, ORANGE)
    pill(slide, "Sandbox", 6.22, 5.13, 1.24, BLUE)
    text(slide, "项目介绍与能力全景", 0.75, 6.48, 3.2, 0.32, 11, MUTED)
    text(slide, "2026 · 技术产品版", 9.70, 6.48, 2.6, 0.32, 11, MUTED, align=PP_ALIGN.RIGHT)


def add_positioning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 1)
    add_header(slide, 2, "一个入口，连接模型、代码与工具", "Project Positioning")
    text(
        slide,
        "Codex Go 将“理解需求—规划任务—调用工具—修改代码—验证结果”\n收敛为可交互、可自动化、可审计的一体化工作流。",
        0.68,
        1.67,
        11.9,
        0.9,
        18,
        TEXT,
        False,
        line_spacing=1.15,
    )
    metrics = [
        ("49", "一级功能模块", CYAN),
        ("1,455", "Go 源文件", GREEN),
        ("544", "测试文件", ORANGE),
        ("3", "主流桌面平台", BLUE),
    ]
    for i, (num, label, accent) in enumerate(metrics):
        x = 0.68 + i * 3.05
        rect(slide, x, 2.84, 2.72, 1.4, PANEL, GRID)
        text(slide, num, x + 0.22, 3.06, 2.25, 0.55, 27, accent, True)
        text(slide, label, x + 0.22, 3.68, 2.25, 0.28, 11, MUTED)
    # Value chain.
    stages = [
        ("01", "理解", "上下文 / 文件 / 历史", CYAN),
        ("02", "规划", "任务拆解 / 多 Agent", GREEN),
        ("03", "执行", "Shell / Patch / MCP", ORANGE),
        ("04", "验证", "测试 / Review / 回滚", BLUE),
    ]
    for i, (no, title_value, sub, accent) in enumerate(stages):
        x = 0.75 + i * 3.02
        circle(slide, x, 5.05, 0.52, accent)
        text(slide, no, x, 5.19, 0.52, 0.2, 9, BG, True, align=PP_ALIGN.CENTER)
        text(slide, title_value, x + 0.72, 5.04, 1.4, 0.3, 15, WHITE, True)
        text(slide, sub, x + 0.72, 5.46, 1.95, 0.42, 10, MUTED)
        if i < 3:
            line(slide, x + 2.35, 5.30, x + 2.83, 5.30, GRID, 2)
    footer(slide, "数据来自当前工作区代码结构统计（2026-07-18）")


def add_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 3, "从“能写代码”到“能完成工程任务”", "Why Codex Go")
    card(slide, 0.68, 1.72, 3.72, 2.0, "传统 AI 助手", "偏重对话与片段生成\n上下文切换频繁\n执行结果依赖人工拼接", RED, "×")
    card(slide, 4.80, 1.72, 3.72, 2.0, "脚本式自动化", "流程固定、适应性有限\n权限边界难统一\n复杂任务缺少反馈闭环", ORANGE, "!")
    card(slide, 8.92, 1.72, 3.72, 2.0, "Codex Go", "理解项目全局\n自主调用工具并持续验证\n以会话沉淀完整过程", GREEN, "✓")
    text(slide, "核心价值", 0.70, 4.28, 2.0, 0.35, 14, CYAN, True)
    values = [
        ("更快交付", "把查找、修改、测试与总结串成闭环", CYAN),
        ("更好集成", "CLI、App Server、MCP 多种入口复用同一能力", GREEN),
        ("更稳治理", "沙箱、审批、执行策略与审计记录贯穿全程", BLUE),
    ]
    for i, (title_value, body, accent) in enumerate(values):
        x = 0.70 + i * 4.08
        rect(slide, x, 4.82, 3.72, 1.52, PANEL2, GRID)
        rect(slide, x + 0.20, 5.03, 0.08, 0.91, accent, radius=False)
        text(slide, title_value, x + 0.52, 4.98, 2.8, 0.34, 15, WHITE, True)
        text(slide, body, x + 0.52, 5.45, 2.83, 0.58, 10, TEXT)
    footer(slide)


def add_capabilities(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 1)
    add_header(slide, 4, "覆盖 AI 编程代理的完整能力栈", "Capability Map")
    caps = [
        ("交互体验", "终端 TUI\nMarkdown / Diff / 搜索\n会话恢复与分支", CYAN, "UI"),
        ("智能执行", "Agent Loop\n计划与工具编排\n自动压缩上下文", GREEN, "AI"),
        ("工程工具", "Shell / Apply Patch\n文件搜索 / Review\n并发统一执行", ORANGE, "TO"),
        ("开放连接", "MCP Client / Server\n插件与 Marketplace\nApp Server / JSON-RPC", BLUE, "IO"),
        ("安全治理", "Sandbox Profile\n审批策略 / Exec Policy\n网络访问约束", RED, "SE"),
        ("运行保障", "登录与多 Provider\nRollout / Telemetry\nDoctor / Completion", LIME, "OP"),
    ]
    for i, (title_value, body, accent, badge) in enumerate(caps):
        row, col = divmod(i, 3)
        card(slide, 0.68 + col * 4.08, 1.75 + row * 2.27, 3.72, 1.86, title_value, body, accent, badge)
    footer(slide, "能力依据 README、目录结构与 parity 清单归纳")


def add_workflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 5, "一次任务，形成可验证的智能闭环", "Agent Workflow")
    steps = [
        ("01", "接收目标", "自然语言 / CLI 参数\nIDE 或服务端请求", CYAN),
        ("02", "构建上下文", "目录、文件、Git Diff\n历史会话与记忆", BLUE),
        ("03", "规划与决策", "任务拆解、模型推理\nSkill / Agent 选择", GREEN),
        ("04", "调用工具", "Shell、补丁、MCP\n外部服务与插件", ORANGE),
        ("05", "验证与沉淀", "测试、Review、结果总结\nRollout 与会话持久化", LIME),
    ]
    for i, (no, title_value, body, accent) in enumerate(steps):
        x = 0.64 + i * 2.52
        circle(slide, x + 0.77, 1.92, 0.72, accent)
        text(slide, no, x + 0.77, 2.13, 0.72, 0.23, 11, BG, True, align=PP_ALIGN.CENTER)
        if i < 4:
            line(slide, x + 1.58, 2.28, x + 2.42, 2.28, GRID, 2.2)
        rect(slide, x, 3.08, 2.25, 2.33, PANEL, GRID)
        text(slide, title_value, x + 0.20, 3.35, 1.85, 0.36, 14, WHITE, True, align=PP_ALIGN.CENTER)
        text(slide, body, x + 0.20, 3.94, 1.85, 0.82, 10, TEXT, align=PP_ALIGN.CENTER, line_spacing=1.15)
        pill(slide, ["PROMPT", "CONTEXT", "PLAN", "ACTION", "PROOF"][i], x + 0.45, 5.72, 1.35, accent)
    text(slide, "人在回路：高风险动作可触发审批；策略决定“何时询问、允许什么、记录什么”。", 1.02, 6.43, 11.1, 0.38, 12, MUTED, align=PP_ALIGN.CENTER)
    footer(slide)


def add_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 1)
    add_header(slide, 6, "分层架构：核心复用，入口自由组合", "Technical Architecture")
    layers = [
        ("体验与接入层", ["TUI", "CLI / Exec", "App Server", "MCP Server"], CYAN),
        ("编排与会话层", ["App Runtime", "Agent / Turn", "Session", "Review"], GREEN),
        ("能力与工具层", ["Shell", "Apply Patch", "MCP Tools", "Skill / Plugin"], ORANGE),
        ("基础与治理层", ["Config / Auth", "Model Provider", "Sandbox", "Telemetry"], BLUE),
    ]
    for r, (name, items, accent) in enumerate(layers):
        y = 1.72 + r * 1.22
        rect(slide, 0.68, y, 2.1, 0.88, accent, None)
        text(slide, name, 0.88, y + 0.26, 1.7, 0.28, 13, BG, True, align=PP_ALIGN.CENTER)
        for c, item in enumerate(items):
            x = 3.10 + c * 2.38
            rect(slide, x, y, 2.08, 0.88, PANEL, GRID)
            text(slide, item, x + 0.10, y + 0.26, 1.88, 0.28, 11, WHITE, True, align=PP_ALIGN.CENTER)
        if r < 3:
            line(slide, 1.73, y + 0.9, 1.73, y + 1.18, accent, 2)
    rect(slide, 0.68, 6.66, 11.96, 0.35, PANEL2, GRID)
    rich_text(
        slide,
        [
            {"text": "设计关键词  ", "color": MUTED, "bold": True, "size": 10},
            {"text": "模块化", "color": CYAN, "bold": True, "size": 10},
            {"text": "  ·  ", "color": MUTED, "size": 10},
            {"text": "协议兼容", "color": GREEN, "bold": True, "size": 10},
            {"text": "  ·  ", "color": MUTED, "size": 10},
            {"text": "可替换 Provider", "color": ORANGE, "bold": True, "size": 10},
            {"text": "  ·  ", "color": MUTED, "size": 10},
            {"text": "跨平台治理", "color": BLUE, "bold": True, "size": 10},
        ],
        3.56,
        6.72,
        6.5,
        0.2,
        align=PP_ALIGN.CENTER,
    )
    footer(slide)


def terminal_window(slide, x, y, w, h):
    rect(slide, x, y, w, h, "071019", GRID)
    rect(slide, x, y, w, 0.48, PANEL2, None, radius=True)
    for i, c in enumerate([RED, ORANGE, GREEN]):
        circle(slide, x + 0.22 + i * 0.26, y + 0.16, 0.12, c)
    text(slide, "codex_go — interactive", x + 0.90, y + 0.14, w - 1.8, 0.18, 9, MUTED, font=MONO)
    text(slide, "● connected", x + w - 1.28, y + 0.14, 1.0, 0.18, 8, GREEN, True, font=MONO, align=PP_ALIGN.RIGHT)


def add_experience(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 7, "交互与自动化，两种体验共用一颗内核", "Product Experience")
    terminal_window(slide, 0.68, 1.70, 7.45, 4.88)
    code = [
        ("› 请分析登录模块并补充边界测试", WHITE),
        ("", TEXT),
        ("• Scanning auth/ and related tests…", MUTED),
        ("• Plan", CYAN),
        ("  1. 定位凭据读取与刷新逻辑", TEXT),
        ("  2. 补充过期、缺失与并发场景", TEXT),
        ("  3. 运行测试并审查 diff", TEXT),
        ("", TEXT),
        ("✓ Patched 2 files   +86  -3", GREEN),
        ("✓ go test ./auth    PASS", GREEN),
        ("", TEXT),
        ("› 已完成：新增 6 个边界用例，无回归。", LIME),
    ]
    for i, (s, c) in enumerate(code):
        text(slide, s, 0.98, 2.40 + i * 0.31, 6.75, 0.24, 10.5, c, i in [0, 3, 8, 9, 11], font=MONO)
    text(slide, "面向开发者", 8.65, 1.78, 2.6, 0.34, 13, CYAN, True)
    card(slide, 8.62, 2.18, 3.72, 1.34, "交互式 TUI", "适合探索、协作与复杂任务\n可视化计划、Diff、审批和状态", CYAN)
    card(slide, 8.62, 3.80, 3.72, 1.34, "非交互 Exec", "适合脚本、CI 与批处理\n稳定输出，便于机器消费", GREEN)
    card(slide, 8.62, 5.42, 3.72, 1.16, "服务化接入", "App Server / JSON-RPC / MCP", ORANGE)
    footer(slide, "界面为能力示意图，非实际运行截图")


def add_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 1)
    add_header(slide, 8, "安全不是附加项，而是执行链路的一部分", "Safety & Governance")
    # Concentric trust zones.
    circle(slide, 0.80, 1.70, 4.75, PANEL)
    circle(slide, 1.36, 2.26, 3.63, BG2, BLUE)
    circle(slide, 2.02, 2.92, 2.31, PANEL2, CYAN)
    circle(slide, 2.65, 3.55, 1.05, LIME)
    text(slide, "任务", 2.65, 3.90, 1.05, 0.24, 14, BG, True, align=PP_ALIGN.CENTER)
    text(slide, "工具权限", 2.05, 3.12, 2.25, 0.25, 11, WHITE, True, align=PP_ALIGN.CENTER)
    text(slide, "沙箱与执行策略", 1.42, 2.45, 3.5, 0.25, 11, BLUE, True, align=PP_ALIGN.CENTER)
    text(slide, "审批 · 网络约束 · 审计记录", 0.92, 1.92, 4.5, 0.25, 11, MUTED, True, align=PP_ALIGN.CENTER)
    controls = [
        ("Sandbox Profile", "限制文件、进程和系统能力边界", CYAN),
        ("Approval Policy", "高风险动作按策略请求人工确认", GREEN),
        ("Exec Policy", "命令级规则定义允许、询问或拒绝", ORANGE),
        ("Session / Rollout", "保留过程与结果，支持追踪和恢复", BLUE),
        ("Network Control", "对外部访问施加可配置约束", RED),
    ]
    for i, (title_value, body, accent) in enumerate(controls):
        y = 1.70 + i * 1.03
        rect(slide, 6.20, y, 6.10, 0.78, PANEL, GRID)
        circle(slide, 6.43, y + 0.20, 0.38, accent)
        text(slide, str(i + 1), 6.43, y + 0.30, 0.38, 0.16, 8, BG, True, align=PP_ALIGN.CENTER)
        text(slide, title_value, 7.02, y + 0.15, 1.78, 0.24, 12, WHITE, True)
        text(slide, body, 8.82, y + 0.15, 3.16, 0.42, 10, TEXT)
    footer(slide)


def add_ecosystem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 9, "开放生态：从命令行延伸到团队基础设施", "Integration Ecosystem")
    # Hub and spokes.
    circle(slide, 5.18, 2.20, 3.0, PANEL, CYAN)
    circle(slide, 6.02, 3.04, 1.32, LIME)
    text(slide, "CODEX\nGO", 6.02, 3.34, 1.32, 0.58, 17, BG, True, align=PP_ALIGN.CENTER)
    nodes = [
        (0.80, 1.85, 2.35, "OpenAI / Provider", "模型与认证", CYAN),
        (0.80, 4.70, 2.35, "CI / Automation", "非交互批处理", GREEN),
        (10.18, 1.85, 2.35, "MCP Servers", "工具与企业系统", ORANGE),
        (10.18, 4.70, 2.35, "IDE / App", "App Server 接入", BLUE),
        (4.85, 5.92, 3.55, "Plugin Marketplace", "能力分发与复用", RED),
    ]
    for x, y, w, title_value, sub, accent in nodes:
        rect(slide, x, y, w, 1.0, PANEL, accent)
        text(slide, title_value, x + 0.16, y + 0.18, w - 0.32, 0.28, 12, WHITE, True, align=PP_ALIGN.CENTER)
        text(slide, sub, x + 0.16, y + 0.57, w - 0.32, 0.22, 9, MUTED, align=PP_ALIGN.CENTER)
        line(slide, x + w / 2, y + 0.5, 6.68, 3.70, accent, 1.4)
    pill(slide, "Windows", 4.50, 1.50, 1.25, CYAN)
    pill(slide, "Linux", 6.03, 1.50, 1.05, GREEN)
    pill(slide, "macOS", 7.36, 1.50, 1.08, ORANGE)
    footer(slide)


def add_progress(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 1)
    add_header(slide, 10, "以协议兼容为基线，持续同步演进", "Project Readiness")
    text(slide, "当前对齐基线", 0.70, 1.72, 2.2, 0.28, 11, MUTED, True)
    rich_text(
        slide,
        [
            {"text": "rust-v0.145.0-alpha.20", "color": WHITE, "bold": True, "size": 18},
            {"text": "  /  ", "color": MUTED, "size": 14},
            {"text": "315195492c", "color": CYAN, "bold": True, "size": 14, "font": MONO},
        ],
        0.70,
        2.08,
        5.0,
        0.4,
    )
    # Donut-like status via segmented bar.
    text(slide, "Parity 清单", 7.75, 1.72, 2.0, 0.28, 11, MUTED, True)
    rect(slide, 7.75, 2.16, 4.48, 0.48, PANEL, None)
    rect(slide, 7.75, 2.16, 3.13, 0.48, GREEN, None, radius=False)
    rect(slide, 10.88, 2.16, 0.90, 0.48, ORANGE, None, radius=False)
    rect(slide, 11.78, 2.16, 0.45, 0.48, BLUE, None, radius=False)
    text(slide, "7 DONE", 7.78, 2.29, 1.7, 0.18, 9, BG, True)
    text(slide, "2 PARTIAL", 10.72, 2.29, 1.0, 0.18, 8, BG, True)
    text(slide, "1 DIFF", 11.78, 2.29, 0.45, 0.18, 7, BG, True, align=PP_ALIGN.CENTER)
    milestones = [
        ("已覆盖", "App Server、Agent 配置、自动压缩、统一执行、Skill 选择、环境状态、TUI 会话历史", GREEN),
        ("持续完善", "外部 Agent 记忆导入语义；macOS 原生沙箱 / root 场景验证", ORANGE),
        ("有意差异", "保留既有 Realtime 能力，外部协议行为保持不变", BLUE),
    ]
    for i, (tag_value, body, accent) in enumerate(milestones):
        y = 3.12 + i * 1.08
        pill(slide, tag_value, 0.72, y, 1.25, accent)
        text(slide, body, 2.22, y - 0.01, 9.85, 0.46, 11, TEXT)
        line(slide, 0.72, y + 0.62, 12.20, y + 0.62, GRID, 0.8)
    text(slide, "工程策略：先锁定协议与行为，再通过测试、发布门禁和跨平台验证逐步收敛。", 0.72, 6.50, 11.5, 0.38, 12, LIME, True, align=PP_ALIGN.CENTER)
    footer(slide, "状态来自 parity.json（generatedAt: 2026-07-18）")


def add_scenarios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, 11, "面向个人、团队与平台的多元场景", "Use Cases")
    scenarios = [
        ("01", "日常研发", "读懂陌生代码\n实现需求与修复缺陷\n生成测试与文档", CYAN),
        ("02", "代码质量", "Review Git Diff\n识别风险与回归\n自动执行验证", GREEN),
        ("03", "研发自动化", "批量重构与迁移\nCI 中运行 Agent\n生成结构化结果", ORANGE),
        ("04", "平台集成", "嵌入 IDE / 工作台\n连接内部工具\n构建专属插件生态", BLUE),
    ]
    for i, (no, title_value, body, accent) in enumerate(scenarios):
        x = 0.68 + i * 3.05
        rect(slide, x, 1.78, 2.72, 3.52, PANEL, GRID)
        circle(slide, x + 0.22, 2.03, 0.52, accent)
        text(slide, no, x + 0.22, 2.17, 0.52, 0.2, 9, BG, True, align=PP_ALIGN.CENTER)
        text(slide, title_value, x + 0.22, 2.82, 2.25, 0.40, 17, WHITE, True)
        text(slide, body, x + 0.22, 3.55, 2.25, 1.08, 11, TEXT, line_spacing=1.2)
        rect(slide, x + 0.22, 4.90, 2.25, 0.06, accent, radius=False)
    rect(slide, 0.68, 5.74, 11.96, 0.86, PANEL2, GRID)
    text(slide, "共同底座", 0.92, 5.98, 1.1, 0.25, 11, MUTED, True)
    rich_text(
        slide,
        [
            {"text": "可配置", "color": CYAN, "bold": True, "size": 13},
            {"text": "   ×   ", "color": MUTED, "size": 11},
            {"text": "可扩展", "color": GREEN, "bold": True, "size": 13},
            {"text": "   ×   ", "color": MUTED, "size": 11},
            {"text": "可治理", "color": ORANGE, "bold": True, "size": 13},
            {"text": "   ×   ", "color": MUTED, "size": 11},
            {"text": "可追踪", "color": BLUE, "bold": True, "size": 13},
        ],
        4.02,
        5.97,
        5.1,
        0.28,
        align=PP_ALIGN.CENTER,
    )
    footer(slide)


def add_cta(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, 1)
    tag(slide, "GET STARTED", 0.72, 0.70, 1.52, LIME)
    text(slide, "让 AI 真正进入\n你的工程工作流", 0.72, 1.46, 6.0, 1.42, 31, WHITE, True, line_spacing=1.05)
    text(slide, "从一次命令开始，体验理解、执行与验证的完整闭环。", 0.75, 3.18, 5.9, 0.60, 15, TEXT)
    terminal_window(slide, 6.62, 1.12, 5.72, 4.64)
    commands = [
        ("# 构建", MUTED),
        ("go build -o ./bin/codex ./cmd/codex", CYAN),
        ("", TEXT),
        ("# 交互式使用", MUTED),
        ('./bin/codex -- "分析当前项目并提出优化建议"', GREEN),
        ("", TEXT),
        ("# 非交互执行", MUTED),
        ('./bin/codex exec "总结本次 Git 变更"', ORANGE),
    ]
    for i, (cmd, c) in enumerate(commands):
        text(slide, cmd, 6.92, 1.92 + i * 0.42, 5.05, 0.28, 10, c, c != MUTED, font=MONO)
    pill(slide, "Go 1.26.2+", 0.75, 4.40, 1.35, CYAN)
    pill(slide, "Windows", 2.27, 4.40, 1.20, BLUE)
    pill(slide, "Linux", 3.64, 4.40, 1.03, GREEN)
    pill(slide, "macOS", 4.84, 4.40, 1.04, ORANGE)
    text(slide, "Codex Go", 0.75, 6.18, 2.5, 0.44, 22, LIME, True)
    text(slide, "Build faster. Integrate deeper. Govern safely.", 0.75, 6.67, 5.7, 0.28, 11, MUTED)
    text(slide, "THANK YOU", 10.20, 6.45, 2.10, 0.35, 13, CYAN, True, align=PP_ALIGN.RIGHT)


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "Codex Go 项目宣传介绍"
    prs.core_properties.subject = "跨平台 AI 编程代理项目介绍"
    prs.core_properties.author = "Codex Go Project"
    prs.core_properties.keywords = "Codex Go, AI 编程代理, Go, MCP, TUI, App Server"

    add_cover(prs)
    add_positioning(prs)
    add_value(prs)
    add_capabilities(prs)
    add_workflow(prs)
    add_architecture(prs)
    add_experience(prs)
    add_security(prs)
    add_ecosystem(prs)
    add_progress(prs)
    add_scenarios(prs)
    add_cta(prs)

    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    build()
